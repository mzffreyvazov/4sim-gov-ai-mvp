# task_d/main.py

import os
import json
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.responses import FileResponse
import fitz  # PyMuPDF
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from datetime import datetime


# Load env vars
load_dotenv()
GOOGLE_API_KEY = ""

app = FastAPI(title="Task D: Slide Generation")

def extract_json_from_response(text: str) -> str:
    """Extract JSON content from markdown code blocks or plain text."""
    text = text.strip()
    
    # Check if it's wrapped in markdown code blocks
    if text.startswith("```json") and text.endswith("```"):
        # Extract content between ```json and ```
        lines = text.split('\n')
        json_lines = []
        in_json = False
        for line in lines:
            if line.strip() == "```json":
                in_json = True
                continue
            elif line.strip() == "```" and in_json:
                break
            elif in_json:
                json_lines.append(line)
        return '\n'.join(json_lines)
    elif text.startswith("```") and text.endswith("```"):
        # Extract content between ``` and ```
        return text[3:-3].strip()
    else:
        return text

def add_title_slide(prs, presentation_data):
    """Add title slide using template approach"""
    # Use the first slide if it exists, otherwise add new one
    if len(prs.slides) > 0:
        slide = prs.slides[0]
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Get current date
    current_date = datetime.now().strftime("%d/%m/%Y")
    
    # Find and update title placeholder
    title_shape = slide.shapes.title
    if title_shape and title_shape.text_frame:
        p = title_shape.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].text = presentation_data['title']
        else:
            p.text = presentation_data['title']
        
        # Format title
        for run in p.runs:
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.name = "Arial"
    
    # Find placeholder for date and update it
    for placeholder in slide.placeholders:
        if placeholder.text and ("Tarix" in placeholder.text or placeholder.text.strip() == ""):
            text_frame = placeholder.text_frame
            p = text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = f"Tarix: {current_date}"
            else:
                p.text = f"Tarix: {current_date}"
            
            # Format date
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.name = "Arial"
            break

def add_intro_slide(prs, presentation_data):
    """Add introduction slide with project content and goals"""
    # Add slide using appropriate layout
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    
    vertical_margin = Inches(0.2)
    
    # Look for existing placeholders or shapes to position content relative to them
    project_content_added = False
    goals_added = False
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            
            # Handle project content placeholder
            if "Layihənin məzmunu" in text or "Project Content" in text:
                if presentation_data.get('project_content'):
                    # Create new text box below the placeholder
                    new_text_box_width = shape.width
                    left = shape.left
                    top = shape.top + shape.height + vertical_margin
                    
                    txBox = slide.shapes.add_textbox(left, top, new_text_box_width, Inches(3))
                    tf = txBox.text_frame
                    
                    # Configure text frame
                    tf.word_wrap = True
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0
                    
                    # Add content
                    p_content = tf.paragraphs[0]
                    p_content.clear()
                    run_content = p_content.add_run()
                    run_content.text = presentation_data['project_content']
                    run_content.font.size = Pt(17)
                    run_content.font.bold = False
                    run_content.font.name = "Arial"
                    
                    project_content_added = True
            
            # Handle goals placeholder
            elif "Məqsəd" in text or "Goals" in text:
                if presentation_data.get('goals'):
                    # Create new text box below the placeholder
                    new_text_box_width = shape.width
                    left = shape.left
                    top = shape.top + shape.height + vertical_margin
                    
                    txBox = slide.shapes.add_textbox(left, top, new_text_box_width, Inches(3))
                    tf = txBox.text_frame
                    
                    # Configure text frame
                    tf.word_wrap = True
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0
                    
                    # Add content
                    p_content = tf.paragraphs[0]
                    p_content.clear()
                    run_content = p_content.add_run()
                    run_content.text = presentation_data['goals']
                    run_content.font.size = Pt(17)
                    run_content.font.bold = False
                    run_content.font.name = "Arial"
                    
                    goals_added = True
    
    # Fallback: if no placeholders found, add content manually
    if not project_content_added and presentation_data.get('project_content'):
        pc_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.29), Inches(4.44), Inches(2))
        tf_pc = pc_box.text_frame
        p_pc = tf_pc.paragraphs[0]
        p_pc.text = presentation_data['project_content']
        p_pc.font.bold = True
        p_pc.font.size = Pt(17)
        p_pc.font.name = "Arial"
    
    if not goals_added and presentation_data.get('goals'):
        goal_box = slide.shapes.add_textbox(Inches(7.11), Inches(1.29), Inches(4.44), Inches(2))
        tf_goal = goal_box.text_frame
        p_goal = tf_goal.paragraphs[0]
        p_goal.text = presentation_data['goals']
        p_goal.font.bold = True
        p_goal.font.size = Pt(17)
        p_goal.font.name = "Arial"

def add_main_slide(prs, slide_data):
    """Add main content slide using template approach"""
    # Use layout 12 or fallback to available layout
    layout_index = min(12, len(prs.slide_layouts) - 1)
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    
    # Set slide title
    if slide.shapes.title:
        slide.shapes.title.text = slide_data['title']
        # Format title
        if slide.shapes.title.text_frame:
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(25)
                    run.font.bold = True
                    run.font.name = "Arial"
    
    # Find empty placeholders and populate them with punkt content
    current_point_idx = 1
    
    for shape in slide.shapes:
        if shape.is_placeholder and current_point_idx <= 4:
            current_text_in_shape = shape.text.strip()
            
            # If placeholder is empty, populate it
            if current_text_in_shape == "":
                punkt_key = f'punkt_{chr(96 + current_point_idx)}_content'  # punkt_a_content, punkt_b_content, etc.
                point_content = slide_data.get(punkt_key, '')
                
                if point_content:
                    tf = shape.text_frame
                    tf.clear()
                    
                    p = tf.add_paragraph()
                    p.text = point_content
                    p.font.size = Pt(17)
                    p.font.name = "Arial"
                else:
                    shape.text_frame.clear()
                
                current_point_idx += 1
    
    # If no placeholders were found, create manual layout
    if current_point_idx == 1:
        # Create manual punkt layout
        punkts = [
            slide_data.get("punkt_a_content", "Punkt A"),
            slide_data.get("punkt_b_content", "Punkt B"),
            slide_data.get("punkt_c_content", "Punkt C"),
            slide_data.get("punkt_d_content", "Punkt D")
        ]
        
        x_start = 0.61
        for idx, content in enumerate(punkts):
            # Circle label
            cx = Inches(x_start + idx * 2.86 + 1.27)
            cy = Inches(2.51)
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                cx, cy,
                Inches(0.31), Inches(0.31)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(0x8B, 0x77, 0x35)
            lbl = circle.text_frame.paragraphs[0]
            lbl.text = chr(65 + idx)  # A, B, C, D
            lbl.font.size = Pt(16)
            lbl.font.bold = True
            lbl.font.name = "Arial"
            
            # Header box
            header = slide.shapes.add_textbox(
                Inches(x_start + idx * 2.86), Inches(2.76),
                Inches(2.86), Inches(1.23)
            )
            th = header.text_frame
            ph = th.paragraphs[0]
            ph.text = f"Punkt {chr(65 + idx)}"  # Punkt A, Punkt B, etc.
            ph.font.bold = True
            ph.font.size = Pt(16)
            ph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            ph.font.name = "Arial"
            header.fill.solid()
            header.fill.fore_color.rgb = RGBColor(0x8B, 0x77, 0x35)
            
            # Content box
            cb = slide.shapes.add_textbox(
                Inches(x_start + idx * 2.86), Inches(4.16),
                Inches(2.86), Inches(1.08)
            )
            tc = cb.text_frame
            pc = tc.paragraphs[0]
            pc.text = content
            pc.font.size = Pt(14)
            pc.font.name = "Arial"
            pc.alignment = PP_ALIGN.LEFT

def add_recommendation_slide(prs, presentation_data):
    """Add next steps/recommendation slide"""
    # Use layout 3 or fallback
    layout_index = min(3, len(prs.slide_layouts) - 1)
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    
    # Set title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Növbəti addımlar"
        # Format title
        if title_shape.text_frame:
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(25)
                    run.font.bold = True
                    run.font.name = "Arial"
    
    # Find empty placeholder for content
    target_tf = None
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        paragraphs = shape.text_frame.paragraphs
        if not paragraphs:
            continue
        
        first_text = paragraphs[0].text.strip()
        
        # Find empty placeholder
        if first_text == "":
            target_tf = shape.text_frame
            break
    
    # If found placeholder, populate it
    if target_tf:
        target_tf.clear()
        
        next_steps = presentation_data.get("next_steps", [])
        for i, step in enumerate(next_steps):
            if step:
                p = target_tf.add_paragraph()
                p.text = f"• {step}"
                p.level = 0
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.font.size = Pt(18)
                p.font.name = "Arial"
    else:
        # Fallback: create manual text box
        next_steps_text = "\n".join([f"• {step}" for step in presentation_data.get("next_steps", [])])
        
        ns_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.29), Inches(10), Inches(3))
        tf_ns = ns_box.text_frame
        p_ns = tf_ns.paragraphs[0]
        p_ns.text = next_steps_text
        p_ns.font.size = Pt(18)
        p_ns.font.name = "Arial"

def delete_slide(prs, slide_index):
    """Delete slide at given index"""
    if slide_index < len(prs.slides._sldIdLst):
        slide_id = prs.slides._sldIdLst[slide_index]
        prs.slides._sldIdLst.remove(slide_id)

@app.post("/parse_doc/")
async def parse_doc(file: UploadFile = File(...)):
    ext = file.filename.rsplit(".", 1)[-1].lower()
    data = await file.read()
    doc = fitz.open(stream=data, filetype=ext)
    text = ""
    images = []
    for page in doc:
        text += page.get_text()
        for img_meta in page.get_images(full=True):
            xref = img_meta[0]
            pix = fitz.Pixmap(doc, xref)
            buf = pix.tobytes("png")
            images.append(buf)
    return {"text": text, "image_count": len(images)}


@app.post("/make_slides/")
async def make_slides(
    doc: UploadFile = File(...),
    prompt: str = Form(...),
    slide_count: int = Form(...)
):
    import logging
    logging.basicConfig(level=logging.INFO)
    logger = logging.getLogger(__name__)

    # A) Parse the document
    parsed = await parse_doc(doc)
    doc_text = parsed["text"]

    # B) Initialize Gemini with API key
    llm = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0.0,
        api_key=GOOGLE_API_KEY
    )

    # C) FIRST REQUEST: Generate all slide content in one go
    content_generation_prompt = PromptTemplate(
        input_variables=["prompt", "text", "slide_count"],
        template=(
            "User instruction: {prompt}\n\n"
            "Document content:\n{text}\n\n"
            "Based on the document and user instruction, generate ALL presentation content for {slide_count} main content slides.\n"
            "Write everything in Azerbaijani language.\n\n"
            "Provide the following content:\n"
            "1. TITLE: A presentation title\n"
            "2. PROJECT_CONTENT: 1-2 sentences summarizing the document\n"
            "3. GOALS: Project goals summary\n"
            "4. SLIDES: For each of the {slide_count} slides, provide:\n"
            "   - Title of the slide\n"
            "   - Content for punkt A (2-3 sentences)\n"
            "   - Content for punkt B (2-3 sentences)\n"
            "   - Content for punkt C (2-3 sentences)\n"
            "   - Content for punkt D (2-3 sentences)\n"
            "5. NEXT_STEPS: Provide exactly 5 bullet points for next steps\n\n"
            "Write in a clear, structured format with headings."
        )
    )
    
    content_chain = content_generation_prompt | llm | StrOutputParser()
    all_content = content_chain.invoke({
        "prompt": prompt,
        "text": doc_text,
        "slide_count": slide_count
    })
    
    logger.info(f"Generated content: {all_content}")

    # D) SECOND REQUEST: Parse the content into structured JSON
    json_parsing_prompt = PromptTemplate(
        input_variables=["content", "slide_count"],
        template=(
            "Parse the following content into a JSON structure with this exact format:\n\n"
            "```json\n"
            "{{\n"
            '  "title": "presentation title",\n'
            '  "project_content": "project content summary",\n'
            '  "goals": "project goals",\n'
            '  "slides": [\n'
            '    {{\n'
            '      "title": "slide title",\n'
            '      "punkt_a_content": "content for punkt A",\n'
            '      "punkt_b_content": "content for punkt B",\n'
            '      "punkt_c_content": "content for punkt C",\n'
            '      "punkt_d_content": "content for punkt D"\n'
            '    }}\n'
            '  ],\n'
            '  "next_steps": [\n'
            '    "step 1",\n'
            '    "step 2",\n'
            '    "step 3",\n'
            '    "step 4",\n'
            '    "step 5"\n'
            '  ]\n'
            "}}\n"
            "```\n\n"
            "Content to parse:\n{content}\n\n"
            "Return ONLY the JSON, no other text. Make sure there are exactly {slide_count} slides in the array."
        )
    )
    
    json_chain = json_parsing_prompt | llm | StrOutputParser()
    json_response = json_chain.invoke({
        "content": all_content,
        "slide_count": slide_count
    })
    
    logger.info(f"JSON response: {json_response}")

    # Parse the JSON response
    clean_json = extract_json_from_response(json_response)
    logger.info(f"Cleaned JSON: {clean_json}")
    
    try:
        presentation_data = json.loads(clean_json)
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error: {e}")
        logger.error(f"Raw JSON: {clean_json}")
        # Fallback with dummy data
        presentation_data = {
            "title": "Təqdimat Başlığı",
            "project_content": "Layihə məzmunu",
            "goals": "Layihə məqsədləri",
            "slides": [
                {
                    "title": "Slayd 1",
                    "punkt_a_content": "Punkt A məzmunu",
                    "punkt_b_content": "Punkt B məzmunu", 
                    "punkt_c_content": "Punkt C məzmunu",
                    "punkt_d_content": "Punkt D məzmunu"
                }
            ],
            "next_steps": [
                "Addım 1", "Addım 2", "Addım 3", "Addım 4", "Addım 5"
            ]
        }

    # E) Assemble PPTX using template-based approach
    try:
        # Try to load your existing template first
        prs = Presentation("taskd/assets/template.pptx")
        logger.info("Loaded existing template")
    except Exception as e:
        logger.warning(f"Could not load template: {e}, creating new presentation")
        # Create a new presentation if template fails
        prs = Presentation()

    # 1. Add title slide
    add_title_slide(prs, presentation_data)
    
    # 2. Add introduction slide
    add_intro_slide(prs, presentation_data)
    
    # 3. Add main content slides
    for slide_data in presentation_data["slides"]:
        add_main_slide(prs, slide_data)
    
    # 4. Add next steps slide
    add_recommendation_slide(prs, presentation_data)
    
    # Clean up any extra slides that might have been in the template
    # (Be careful with this - only delete if you're sure about the template structure)
    try:
        # If template had extra slides, you might want to remove them
        while len(prs.slides) > len(presentation_data["slides"]) + 3:  # +3 for title, intro, next steps
            delete_slide(prs, -1)  # Remove last slide
    except Exception as e:
        logger.warning(f"Could not clean up extra slides: {e}")

    # Save the presentation
    output_path = "generated_presentation.pptx"
    prs.save(output_path)

    return FileResponse(
        output_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="presentation.pptx"
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)