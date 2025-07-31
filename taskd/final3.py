import os
import json
import asyncio
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
import fitz
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
import logging
from datetime import datetime
from jinja2 import Environment, FileSystemLoader

# Imports for the robust PPTX generation
from pptx import Presentation
from pptx.util import Inches
from playwright.async_api import async_playwright

# --- 1. SETUP ---
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
app = FastAPI(title="Task D: Slide Generation")
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# --- 2. HELPER FUNCTIONS ---
def extract_json_from_response(text: str) -> str:
    """Extracts a JSON object from a string, handling markdown and other noise."""
    try:
        json_start = text.index('{')
        json_end = text.rindex('}') + 1
        return text[json_start:json_end]
    except ValueError:
        raise ValueError("No valid JSON object found in the LLM response.")

async def parse_doc(file: UploadFile = File(...)) -> str:
    """Parses an uploaded document (PDF, etc.) to extract raw text."""
    try:
        ext = file.filename.rsplit(".", 1)[-1].lower()
        data = await file.read()
        doc = fitz.open(stream=data, filetype=ext)
        text = "".join(page.get_text() for page in doc)
        return text
    except Exception as e:
        logger.error(f"Error parsing document: {e}")
        raise HTTPException(status_code=400, detail="Failed to parse the document.")


# --- 3. HTML GENERATION ---
def generate_html_slides(data: dict, base_dir: str) -> list[str]:
    """Uses Jinja2 to populate HTML templates with JSON data and save them as files."""
    templates_dir = os.path.join(base_dir, "html-templates")
    env = Environment(loader=FileSystemLoader(templates_dir))
    
    # Create the output directory relative to the script's location
    output_dir = os.path.join(base_dir, "output_html")
    os.makedirs(output_dir, exist_ok=True)
    
    generated_files = []
    slide_counter = 1
    presentation_data = data.get("presentation", {})
    assets_path = os.path.join(templates_dir, "assets").replace("\\", "/")

    # Generate Intro Slide
    intro_data = presentation_data.get("intro_slide", {})
    if intro_data:
        template = env.get_template("slide-basliq.html")
        rendered_html = template.render(
            presentation_title=intro_data.get("presentation_title", "N/A"),
            presentation_date=intro_data.get("presentation_date", ""),
            assets_path=assets_path
        )
        file_path = os.path.join(output_dir, f"{slide_counter:02d}_intro.html")
        with open(file_path, "w", encoding="utf-8") as f: f.write(rendered_html)
        generated_files.append(file_path)
        slide_counter += 1

    # Generate Project Content & Goal Slide
    content_goal_data = presentation_data.get("project_content_slide", {})
    if content_goal_data:
        template = env.get_template("slide-giris_and_novbeti.html")
        rendered_html = template.render(
            left_panel_title="Layihənin Məzmunu", right_panel_title="Məqsəd",
            left_panel_content=content_goal_data.get("presentation_overview", ""),
            right_panel_content=content_goal_data.get("presentation_goal", ""),
            page_number=slide_counter, assets_path=assets_path
        )
        file_path = os.path.join(output_dir, f"{slide_counter:02d}_content_goal.html")
        with open(file_path, "w", encoding="utf-8") as f: f.write(rendered_html)
        generated_files.append(file_path)
        slide_counter += 1

    # Generate Main Content Slides (Punkts)
    content_slides = presentation_data.get("content_slides", [])
    if content_slides:
        template = env.get_template("slide-punkts.html")
        for i, slide in enumerate(content_slides):
            rendered_html = template.render(
                general_content_title=slide.get("general_content_title", "Məzmun"),
                contents=slide.get("contents", []), page_number=slide_counter,
                assets_path=assets_path
            )
            file_path = os.path.join(output_dir, f"{slide_counter:02d}_content_{i+1}.html")
            with open(file_path, "w", encoding="utf-8") as f: f.write(rendered_html)
            generated_files.append(file_path)
            slide_counter += 1
            
    # Generate Final Slide (Next Steps)
    final_data = presentation_data.get("final_slide", {})
    if final_data:
        next_steps_list = final_data.get("next_steps", [])
        next_steps_html = "<ul>" + "".join(f"<li>{step}</li>" for step in next_steps_list) + "</ul>"
        template = env.get_template("slide-giris_and_novbeti.html")
        rendered_html = template.render(
            left_panel_title="Növbəti Addımlar", right_panel_title="",
            left_panel_content=next_steps_html, right_panel_content="",
            page_number=slide_counter, assets_path=assets_path
        )
        file_path = os.path.join(output_dir, f"{slide_counter:02d}_final.html")
        with open(file_path, "w", encoding="utf-8") as f: f.write(rendered_html)
        generated_files.append(file_path)

    logger.info(f"Generated {len(generated_files)} HTML files in '{output_dir}'")
    return generated_files


# --- 4. ASYNC PPTX GENERATION ---
async def create_pptx_with_playwright(html_files: list[str], output_pptx_path: str) -> str:
    """Renders HTML files to images using Playwright's async API and assembles them into a PPTX."""
    image_files = []
    screenshot_dir = os.path.join(os.path.dirname(output_pptx_path), "screenshots")
    os.makedirs(screenshot_dir, exist_ok=True)
    logger.info("Starting async Playwright to generate screenshots...")

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page()
        await page.set_viewport_size({"width": 1920, "height": 1080})

        for i, html_file in enumerate(html_files):
            file_uri = f"file:///{os.path.abspath(html_file)}"
            await page.goto(file_uri, wait_until="networkidle")
            image_path = os.path.join(screenshot_dir, f"slide_{i:02d}.png")
            await page.screenshot(path=image_path)
            image_files.append(image_path)

        await browser.close()
    logger.info("Screenshots generated. Assembling PowerPoint presentation...")

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]

    for image_path in image_files:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
    prs.save(output_pptx_path)
    logger.info(f"Presentation saved to {output_pptx_path}")
    return output_pptx_path


# --- 5. FASTAPI ENDPOINT ---
@app.post("/make_slides/")
async def make_slides(
    doc: UploadFile = File(...),
    prompt: str = Form(...),
    slide_count: int = Form(...)
):
    """Orchestrates the full process from document to final .pptx file."""
    # Define base directory for this script
    base_dir = os.path.dirname(__file__)

    # --- Part 1: Generate JSON from Document ---
    try:
        doc_text = await parse_doc(doc)
        llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.1, api_key=GOOGLE_API_KEY)
        
        json_format_string = """{"presentation": {"intro_slide": ..., "project_content_slide": ..., "content_slides": [...], "chart_slides": [], "final_slide": ...}}"""
        generation_prompt = PromptTemplate(
            input_variables=["prompt", "text", "slide_count", "json_format", "current_date"],
            template=(
                "User instruction: {prompt}\n\nDocument content:\n---\n{text}\n---\n\n"
                "**CRITICAL RULES:**\n"
                "1. `content_slides` array MUST contain EXACTLY {slide_count} objects.\n"
                "2. Inside EACH `content_slides` object, the `contents` array MUST contain EXACTLY 4 objects.\n"
                "3. The `presentation_title` in `intro_slide` MUST be a maximum of 4-5 words.\n\n"
                "Return ONLY a single valid JSON object.\n\nJSON Format Example:\n{json_format}"
            )
        )
        generation_chain = generation_prompt | llm | StrOutputParser()
        current_date_str = datetime.now().strftime("%d.%m.%Y")
        
        logger.info("Generating slide content from LLM...")
        llm_response = await generation_chain.ainvoke({
            "prompt": prompt, "text": doc_text, "slide_count": slide_count,
            "json_format": json_format_string, "current_date": current_date_str
        })
        
        clean_json_str = extract_json_from_response(llm_response)
        output_json = json.loads(clean_json_str)
        logger.info("Successfully generated and parsed JSON from LLM.")
    except Exception as e:
        logger.error(f"Failed during LLM/JSON generation: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed during LLM or JSON parsing stage. Error: {e}")

    # --- Part 2: Generate HTML files ---
    try:
        html_files = generate_html_slides(output_json, base_dir)
        if not html_files:
            raise ValueError("HTML slide generation resulted in no files.")
    except Exception as e:
        logger.error(f"Failed during HTML generation stage: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to generate HTML slides. Error: {e}")

    # --- Part 3: Convert HTML to PPTX using ASYNC Playwright ---
    try:
        output_dir = os.path.join(base_dir, "final_presentations")
        os.makedirs(output_dir, exist_ok=True)
        pptx_filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        pptx_path = os.path.join(output_dir, pptx_filename)
        
        final_file_path = await create_pptx_with_playwright(html_files, pptx_path)
        
        return FileResponse(
            path=final_file_path, 
            filename=pptx_filename, 
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        logger.error(f"Failed during async PPTX generation stage: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to generate PPTX file. Error: {e}")