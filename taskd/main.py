import os
import json
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request, Body
from fastapi.responses import JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
import fitz
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
import logging
from datetime import datetime
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
app = FastAPI(title="Task D: Slide Generation")

# Mount static files
app.mount("/static", StaticFiles(directory="static"), name="static")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Global variable to store the current JSON data
current_presentation_data = {}

def extract_json_from_response(text: str):
    try:
        json_start = text.index('{')
        json_end = text.rindex('}') + 1
        return text[json_start:json_end]
    except ValueError:
        raise ValueError("No valid JSON object found in the LLM response.")

async def parse_doc(file: UploadFile = File(...)):
    try:
        ext = file.filename.rsplit(".", 1)[-1].lower()
        data = await file.read()
        doc = fitz.open(stream=data, filetype=ext)
        text = "".join(page.get_text() for page in doc)
        return {"text": text}
    except Exception as e:
        logger.error(f"Error parsing document: {e}")
        raise HTTPException(status_code=400, detail="Failed to parse the document.")


def debug_template_structure(template_path: str):
    try:
        prs = Presentation(template_path)
        logger.info(f"Template has {len(prs.slide_layouts)} slide layouts:")
        
        for i, layout in enumerate(prs.slide_layouts):
            logger.info(f"  Layout {i}: {layout.name} - {len(layout.placeholders)} placeholders")
            for j, placeholder in enumerate(layout.placeholders):
                placeholder_name = getattr(placeholder, 'name', f'unnamed_{j}')
                logger.info(f"    Placeholder {j}: {placeholder_name}")
                
    except Exception as e:
        logger.error(f"Error examining template structure: {e}")

def generate_pptx_presentation(data: dict) -> str:
    """
    Generate a PowerPoint presentation using python-pptx from JSON data.
    Uses the template_final.pptx file with custom slide layouts matching the JSON structure.
    """
    # Get template path
    template_path = os.path.join(os.path.dirname(__file__), "html-templates", "template_final.pptx")
    
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    # Debug template structure (for troubleshooting)
    debug_template_structure(template_path)
    
    # Load the presentation template
    prs = Presentation(template_path)
    
    # Set output directory and filename with timestamp
    output_dir = os.path.join(os.path.dirname(__file__), "output_folder_name")
    os.makedirs(output_dir, exist_ok=True)
    
    current_time = datetime.now()
    filename = f"output_presentation_{current_time.strftime('%Y%m%d')}_{current_time.strftime('%H%M%S')}.pptx"
    output_path = os.path.join(output_dir, filename)
    
    presentation_data = data.get("presentation", {})
    
    # Helper function to find slide layout by name
    def find_layout_by_name(layout_name: str):
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
        return None
    
    # Helper function to find placeholder by name
    def find_placeholder_by_name(slide, placeholder_name: str):
        for shape in slide.placeholders:
            if hasattr(shape, 'name') and shape.name == placeholder_name:
                return shape
        return None
    
    # Helper function to adjust placeholder height based on text length
    def adjust_placeholder_height(placeholder, text: str, placeholder_type: str = "default"):
        """
        Adjust the height of a placeholder based on the length of text content.
        Also ensures appropriate width and font size based on placeholder type.
        """
        if not placeholder or not text:
            return
            
        try:
            from pptx.util import Inches, Pt
            
            # Set dimensions and font size based on placeholder type
            if placeholder_type == "content":
                # For content-1, content-2, etc. placeholders
                fixed_width = Inches(2.80)
                max_height = Inches(2.40)
                font_size = Pt(14)
                char_per_line = 30  # Shorter lines for narrower width
            elif placeholder_type == "overview"  or  placeholder_type == "goal":
                # For presentation_overview placeholder
                fixed_width = Inches(5.55)
                max_height = Inches(4.0)  # Allow more height for overview
                font_size = Pt(16)
                char_per_line = 60  # More characters for wider width
            elif placeholder_type == "chart_summary":
                # For chart summary placeholder (Text Placeholder 2)
                fixed_width = Inches(6.76)
                max_height = Inches(4.0)  # Allow sufficient height for chart summaries
                font_size = Pt(14)
                char_per_line = 70  # More characters for wider width
            else:
                # Default for other placeholders
                fixed_width = Inches(6.5)
                max_height = Inches(4.0)
                font_size = Pt(16)
                char_per_line = 50
            
            # Set fixed width
            placeholder.width = fixed_width
            
            # Calculate approximate height needed based on text length
            base_height = Inches(0.5)  # Minimum height
            lines_needed = max(1, len(text) // char_per_line + text.count('\n') + 1)
            
            # Calculate new height (roughly 0.3 inches per line)
            new_height = base_height + Inches(0.3 * (lines_needed - 1))
            
            # Apply maximum height limit
            new_height = min(new_height, max_height)
            
            # Apply the new height
            placeholder.height = new_height
            
            # Set font size if text frame exists
            if placeholder.has_text_frame and placeholder.text_frame.paragraphs:
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = font_size
            
            logger.info(f"Adjusted {placeholder_type} placeholder to width: {fixed_width}, height: {new_height}, font: {font_size} for {len(text)} characters, {lines_needed} lines")
        except Exception as e:
            logger.error(f"Error adjusting placeholder height: {e}")

    def configure_chart_legend(chart, chart_type: str):
        """
        Configure chart legend settings for better visibility.
        """
        try:
            # Enable legend for all charts
            chart.has_legend = True
            
            # Position legend based on chart type
            if chart_type == "pie":
                # For pie charts, place legend on the right side
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                
                # Enable data labels with percentages for pie charts
                try:
                    # Get the first (and usually only) series in the chart
                    if len(chart.series) > 0:
                        series = chart.series[0]
                        # Enable data labels
                        series.has_data_labels = True
                        data_labels = series.data_labels
                        
                        # Configure data labels to show percentages
                        data_labels.show_percentage = True
                        data_labels.show_value = False  # Hide raw values, show only percentages
                        data_labels.show_category_name = False  # Category names are in legend
                        
                        # Set data label font size
                        if hasattr(data_labels, 'font'):
                            data_labels.font.size = Pt(8)
                            data_labels.font.bold = True
                        
                        logger.info("Enabled percentage data labels for pie chart")
                except Exception as label_error:
                    logger.warning(f"Could not configure pie chart data labels: {label_error}")
                    
            else:
                # For other charts, place legend at the bottom
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            
            # Don't include legend in layout calculations to prevent chart size reduction
            chart.legend.include_in_layout = False
            
            # Configure legend font (if available)
            if hasattr(chart.legend, 'font'):
                chart.legend.font.size = Pt(8)  # Set readable font size
                
            logger.info(f"Configured legend for {chart_type} chart")
            
        except Exception as e:
            logger.warning(f"Could not fully configure legend: {e}")

    def create_chart_in_placeholder(placeholder, chart_data):
        """
        Create a chart from chart_data and insert it into the placeholder.
        chart_data should follow the format from chart_json_output_format.json
        """
        try:
            if not chart_data or not placeholder:
                logger.warning("No chart data or placeholder provided")
                return None
                
            # Check if this placeholder can accept charts
            if not hasattr(placeholder, 'insert_chart'):
                logger.warning(f"Placeholder {getattr(placeholder, 'name', 'unnamed')} cannot accept charts")
                return None
                
            # Extract chart information
            chart_type = chart_data.get("type", "").lower()
            title = chart_data.get("title", "Chart")
            x_values = chart_data.get("x", [])
            y_values = chart_data.get("y", [])
            labels = chart_data.get("labels", [])
            sizes = chart_data.get("sizes", [])
            
            # Create ChartData object
            chart_data_obj = ChartData()
            
            # Map chart types to python-pptx chart types
            chart_type_mapping = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "scatter": XL_CHART_TYPE.XY_SCATTER,
                "area": XL_CHART_TYPE.AREA
            }
            
            pptx_chart_type = chart_type_mapping.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Set up chart data based on chart type
            if chart_type == "pie" and labels and sizes:
                # For pie charts, use labels as categories and sizes as values
                chart_data_obj.categories = labels
                chart_data_obj.add_series('Series 1', tuple(sizes))
                
            elif x_values and y_values:
                # For other chart types, use x as categories and y as values
                chart_data_obj.categories = x_values
                chart_data_obj.add_series('Series 1', tuple(y_values))
                
            else:
                # Fallback: create a simple chart with sample data if data is missing
                logger.warning(f"Insufficient chart data for {chart_type}, using sample data")
                chart_data_obj.categories = ['Category 1', 'Category 2', 'Category 3']
                chart_data_obj.add_series('Series 1', (10, 20, 30))
            
            # Insert chart into placeholder
            graphic_frame = placeholder.insert_chart(pptx_chart_type, chart_data_obj)
            chart = graphic_frame.chart
            
            # Set chart title if provided
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
            
            # Configure legend for better visibility
            configure_chart_legend(chart, chart_type)
                
            logger.info(f"Successfully created {chart_type} chart with title '{title}'")
            return chart
            
        except Exception as e:
            logger.error(f"Error creating chart: {e}")
            return None

    def add_chart_to_slide(slide, chart_data, x=None, y=None, width=None, height=None):
        """
        Add a chart directly to slide using shapes.add_chart instead of placeholder.
        This is a fallback when no chart placeholder is available.
        """
        try:
            if not chart_data:
                logger.warning("No chart data provided for direct chart creation")
                return None
                
            # Extract chart information
            chart_type = chart_data.get("type", "").lower()
            title = chart_data.get("title", "Chart")
            x_values = chart_data.get("x", [])
            y_values = chart_data.get("y", [])
            labels = chart_data.get("labels", [])
            sizes = chart_data.get("sizes", [])
            
            # Create ChartData object
            chart_data_obj = ChartData()
            
            # Map chart types to python-pptx chart types
            chart_type_mapping = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "scatter": XL_CHART_TYPE.XY_SCATTER,
                "area": XL_CHART_TYPE.AREA
            }
            
            pptx_chart_type = chart_type_mapping.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Set up chart data based on chart type
            if chart_type == "pie" and labels and sizes:
                # For pie charts, use labels as categories and sizes as values
                chart_data_obj.categories = labels
                chart_data_obj.add_series('Series 1', tuple(sizes))
                
            elif x_values and y_values:
                # For other chart types, use x as categories and y as values
                chart_data_obj.categories = x_values
                chart_data_obj.add_series('Series 1', tuple(y_values))
                
            else:
                # Fallback: create a simple chart with sample data if data is missing
                logger.warning(f"Insufficient chart data for {chart_type}, using sample data")
                chart_data_obj.categories = ['Category 1', 'Category 2', 'Category 3']
                chart_data_obj.add_series('Series 1', (10, 20, 30))
            
            # Set better default position and size if not provided
            # Position chart to match the template specifications
            if x is None:
                x = Inches(5.41)  # Horizontal position from template
            if y is None:
                y = Inches(1.31)  # Vertical position from template
            if width is None:
                width = Inches(7.15)  # Width from template
            if height is None:
                height = Inches(3.71)  # Height from template
            
            # Add chart directly to slide
            graphic_frame = slide.shapes.add_chart(pptx_chart_type, x, y, width, height, chart_data_obj)
            chart = graphic_frame.chart
            
            # Set chart title if provided
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
            
            # Configure legend for better visibility
            configure_chart_legend(chart, chart_type)
                
            logger.info(f"Successfully added {chart_type} chart directly to slide with title '{title}' at position ({x}, {y})")
            return chart
            
        except Exception as e:
            logger.error(f"Error adding chart directly to slide: {e}")
            return None
            
        except Exception as e:
            logger.warning(f"Could not adjust placeholder dimensions: {e}")
    
    try:
        # 1. Generate Intro Slide
        intro_data = presentation_data.get("intro_slide", {})
        if intro_data:
            layout = find_layout_by_name("intro_slide")
            logger.info(f"Looking for layout 'intro_slide', found: {layout is not None}")
            if layout:
                slide = prs.slides.add_slide(layout)
                logger.info(f"Added intro slide, placeholders count: {len(slide.placeholders)}")
                
                # Debug: list all placeholders
                for i, ph in enumerate(slide.placeholders):
                    ph_name = getattr(ph, 'name', f'unnamed_{i}')
                    logger.info(f"  Placeholder {i}: {ph_name}")
                
                # Set presentation title (using "Title 2")
                title_placeholder = find_placeholder_by_name(slide, "Title 2")
                logger.info(f"Found title placeholder: {title_placeholder is not None}")
                if title_placeholder and title_placeholder.has_text_frame:
                    title_placeholder.text = intro_data.get("presentation_title", "")
                    logger.info(f"Set title to: {intro_data.get('presentation_title', '')}")
                
                # Set presentation date (using "Subtitle 1")
                date_placeholder = find_placeholder_by_name(slide, "Subtitle 1")
                logger.info(f"Found date placeholder: {date_placeholder is not None}")
                if date_placeholder and date_placeholder.has_text_frame:
                    date_placeholder.text = intro_data.get("presentation_date", "")
                    logger.info(f"Set date to: {intro_data.get('presentation_date', '')}")
            else:
                logger.warning("intro_slide layout not found")
        
        # 2. Generate Project Content & Goal Slide
        content_goal_data = presentation_data.get("project_content_slide", {})
        if content_goal_data:
            layout = find_layout_by_name("project_content_slide") 
            logger.info(f"Looking for layout 'project_content_slide', found: {layout is not None}")
            if layout:
                slide = prs.slides.add_slide(layout)
                logger.info(f"Added content goal slide, placeholders count: {len(slide.placeholders)}")
                
                # Debug: list all placeholders
                for i, ph in enumerate(slide.placeholders):
                    ph_name = getattr(ph, 'name', f'unnamed_{i}')
                    logger.info(f"  Placeholder {i}: {ph_name}")
                
                # Set presentation overview (using "Title 1")
                overview_placeholder = find_placeholder_by_name(slide, "Title 1")
                logger.info(f"Found overview placeholder: {overview_placeholder is not None}")
                if overview_placeholder and overview_placeholder.has_text_frame:
                    overview_text = content_goal_data.get("presentation_overview", "")
                    overview_placeholder.text = overview_text
                    adjust_placeholder_height(overview_placeholder, overview_text, "overview")
                    logger.info(f"Set overview text length: {len(overview_text)}")
                
                # Set presentation goal (using "Subtitle 2")
                goal_placeholder = find_placeholder_by_name(slide, "Subtitle 2")
                logger.info(f"Found goal placeholder: {goal_placeholder is not None}")
                if goal_placeholder and goal_placeholder.has_text_frame:
                    goal_text = content_goal_data.get("presentation_goal", "")
                    goal_placeholder.text = goal_text
                    adjust_placeholder_height(goal_placeholder, goal_text, "goal")
                    logger.info(f"Set goal text length: {len(goal_text)}")
            else:
                logger.warning("project_content_slide layout not found")
        
        # 3. Generate Content Slides (Punkts)
        content_slides = presentation_data.get("content_slides", [])
        for slide_index, slide_data in enumerate(content_slides):
            layout = find_layout_by_name("content_slides")
            logger.info(f"Looking for layout 'content_slides', found: {layout is not None}")
            if layout:
                slide = prs.slides.add_slide(layout)
                logger.info(f"Added content slide {slide_index + 1}, placeholders count: {len(slide.placeholders)}")
                
                # Debug: list all placeholders
                for i, ph in enumerate(slide.placeholders):
                    ph_name = getattr(ph, 'name', f'unnamed_{i}')
                    logger.info(f"  Placeholder {i}: {ph_name}")
                
                # Set general content title (using "Title 1")
                title_placeholder = find_placeholder_by_name(slide, "Title 1")
                logger.info(f"Found title placeholder: {title_placeholder is not None}")
                if title_placeholder and title_placeholder.has_text_frame:
                    title_placeholder.text = slide_data.get("general_content_title", "")
                    logger.info(f"Set title to: {slide_data.get('general_content_title', '')}")
                
                # Set content points
                contents = slide_data.get("contents", [])
                logger.info(f"Processing {len(contents)} content items")
                
                # Populate content placeholders using actual placeholder names
                placeholder_mapping = [
                    "Text Placeholder 2",  # content-1
                    "Text Placeholder 3",  # content-2
                    "Text Placeholder 4",  # content-3
                    "Text Placeholder 5"   # content-4
                ]
                
                for i, content_item in enumerate(contents[:4]):  # Limit to 4 items
                    if not content_item or i >= len(placeholder_mapping):
                        continue
                    
                    # Find the content field (could be content-1, content-2, etc.) - only use content, not title
                    content_text = ""
                    for key, value in content_item.items():
                        if key != "title" and key.startswith("content"):
                            content_text = value
                            break
                    
                    # Find placeholder using actual name from template
                    placeholder_name = placeholder_mapping[i]
                    content_placeholder = find_placeholder_by_name(slide, placeholder_name)
                    logger.info(f"Looking for placeholder '{placeholder_name}', found: {content_placeholder is not None}")
                    
                    if content_placeholder and content_placeholder.has_text_frame:
                        content_placeholder.text = content_text
                        # Adjust height based on content length with specific content dimensions
                        adjust_placeholder_height(content_placeholder, content_text, "content")
                        logger.info(f"Set content {i+1} to: {content_text[:50]}...")
                    else:
                        logger.warning(f"Could not populate {placeholder_name}")
            else:
                logger.warning("content_slides layout not found")
        
        # 4. Generate Chart Slides
        chart_slides = presentation_data.get("chart_slides", [])
        if chart_slides:
            for chart_slide_data in chart_slides:
                layout = find_layout_by_name("chart_slides")
                logger.info(f"Looking for layout 'chart_slides', found: {layout is not None}")
                if layout:
                    slide = prs.slides.add_slide(layout)
                    logger.info(f"Added chart slide, placeholders count: {len(slide.placeholders)}")
                    
                    # Debug: list all placeholders
                    for i, ph in enumerate(slide.placeholders):
                        ph_name = getattr(ph, 'name', f'unnamed_{i}')
                        logger.info(f"  Chart slide placeholder {i}: {ph_name}")
                    
                    # Set content title (first placeholder - usually Title or Title 1)
                    title_placeholder = None
                    # Try common title placeholder names
                    for title_name in ["Title", "Title 1", "content_title"]:
                        title_placeholder = find_placeholder_by_name(slide, title_name)
                        if title_placeholder:
                            break
                    
                    if not title_placeholder and len(slide.placeholders) > 0:
                        title_placeholder = slide.placeholders[0]  # Fallback to first placeholder
                    
                    if title_placeholder and title_placeholder.has_text_frame:
                        content_title = chart_slide_data.get("content_title", "Chart Title")
                        title_placeholder.text = content_title
                        logger.info(f"Set chart title: {content_title}")
                    
                    # Find chart data and create charts
                    chart_data_list = chart_slide_data.get("chart_data", [])
                    if chart_data_list:
                        # Find suitable placeholder for chart (usually second placeholder or one named chart/content)
                        chart_placeholder = None
                        
                        # Try to find chart placeholder by name
                        for chart_name in ["Chart", "Content Placeholder", "Subtitle", "chart_summary"]:
                            chart_placeholder = find_placeholder_by_name(slide, chart_name)
                            if chart_placeholder:
                                break
                        
                        # If not found by name, try placeholders that can accept charts
                        if not chart_placeholder:
                            for placeholder in slide.placeholders:
                                try:
                                    # Check if this placeholder can accept charts
                                    if hasattr(placeholder, 'insert_chart'):
                                        chart_placeholder = placeholder
                                        break
                                except:
                                    continue
                        
                        # If still not found, use second placeholder as fallback
                        if not chart_placeholder and len(slide.placeholders) > 1:
                            chart_placeholder = slide.placeholders[1]
                        
                        if chart_placeholder and chart_data_list:
                            # Create chart from first chart data entry
                            chart_data = chart_data_list[0] if isinstance(chart_data_list, list) else chart_data_list
                            created_chart = None
                            
                            try:
                                # First try to create chart in placeholder
                                created_chart = create_chart_in_placeholder(chart_placeholder, chart_data)
                                if created_chart:
                                    logger.info("Successfully created chart in placeholder")
                                else:
                                    logger.warning("Placeholder doesn't support charts, trying direct chart creation")
                                    # Fallback: add chart directly to slide with template positioning
                                    created_chart = add_chart_to_slide(slide, chart_data, Inches(5.41), Inches(1.31), Inches(7.15), Inches(3.71))
                                    if created_chart:
                                        logger.info("Successfully created chart directly on slide")
                                    else:
                                        logger.warning("Failed to create chart both in placeholder and directly")
                                        
                            except Exception as e:
                                logger.error(f"Error creating chart: {e}")
                                # Final fallback: add chart directly to slide
                                try:
                                    created_chart = add_chart_to_slide(slide, chart_data, Inches(5.41), Inches(1.31), Inches(7.15), Inches(3.71))
                                    if created_chart:
                                        logger.info("Successfully created chart directly on slide as fallback")
                                    else:
                                        # Ultimate fallback: add chart summary as text
                                        if chart_placeholder and chart_placeholder.has_text_frame:
                                            summary = chart_slide_data.get("chart_summary", "Chart data visualization")
                                            chart_placeholder.text = summary
                                            logger.info("Added chart summary as text fallback")
                                except Exception as e2:
                                    logger.error(f"All chart creation methods failed: {e2}")
                                    # Ultimate fallback: add chart summary as text
                                    if chart_placeholder and chart_placeholder.has_text_frame:
                                        summary = chart_slide_data.get("chart_summary", "Chart data visualization")
                                        chart_placeholder.text = summary
                                        logger.info("Added chart summary as text fallback")
                        
                        else:
                            # No suitable placeholder found, add chart directly to slide
                            chart_data = chart_data_list[0] if isinstance(chart_data_list, list) else chart_data_list
                            try:
                                created_chart = add_chart_to_slide(slide, chart_data, Inches(5.41), Inches(1.31), Inches(7.15), Inches(3.71))
                                if created_chart:
                                    logger.info("Successfully created chart directly on slide (no suitable placeholder)")
                                else:
                                    logger.warning("Failed to create chart directly on slide")
                            except Exception as e:
                                logger.error(f"Error creating chart directly on slide: {e}")
                        
                        # Add chart summary to Text Placeholder 2
                        chart_summary = chart_slide_data.get("chart_summary", "")
                        if chart_summary:
                            # Find Text Placeholder 2 specifically for chart summary
                            summary_placeholder = find_placeholder_by_name(slide, "Text Placeholder 2")
                            if summary_placeholder and summary_placeholder.has_text_frame:
                                summary_placeholder.text = chart_summary
                                adjust_placeholder_height(summary_placeholder, chart_summary, "chart_summary")
                                logger.info(f"Added chart summary to Text Placeholder 2: {chart_summary[:50]}...")
                            else:
                                logger.warning("Text Placeholder 2 not found for chart summary")
                                # Fallback: Look for any remaining text placeholders for summary
                                for placeholder in slide.placeholders:
                                    if (placeholder != title_placeholder and 
                                        placeholder != chart_placeholder and 
                                        placeholder.has_text_frame and 
                                        not placeholder.text.strip()):
                                        placeholder.text = chart_summary
                                        logger.info(f"Added chart summary to fallback placeholder: {chart_summary[:50]}...")
                                        break
                    
                    else:
                        # No chart data provided, add summary text to Text Placeholder 2
                        chart_summary = chart_slide_data.get("chart_summary", "")
                        if chart_summary:
                            summary_placeholder = find_placeholder_by_name(slide, "Text Placeholder 2")
                            if summary_placeholder and summary_placeholder.has_text_frame:
                                summary_placeholder.text = chart_summary
                                adjust_placeholder_height(summary_placeholder, chart_summary, "chart_summary")
                                logger.info(f"Added chart summary to Text Placeholder 2: {chart_summary[:50]}...")
                            else:
                                logger.warning("Text Placeholder 2 not found, trying fallback")
                                # Fallback to any available text placeholder
                                if len(slide.placeholders) > 1:
                                    summary_placeholder = slide.placeholders[1]
                                    if summary_placeholder and summary_placeholder.has_text_frame:
                                        summary_placeholder.text = chart_summary
                                        logger.info(f"Added chart summary text to fallback placeholder: {chart_summary[:50]}...")
                else:
                    logger.warning("chart_slides layout not found")
        
        # 5. Generate Final Slide (Next Steps)
        final_data = presentation_data.get("final_slide", {})
        if final_data:
            layout = find_layout_by_name("final_slide")
            logger.info(f"Looking for layout 'final_slide', found: {layout is not None}")
            if layout:
                slide = prs.slides.add_slide(layout)
                logger.info(f"Added final slide, placeholders count: {len(slide.placeholders)}")
                
                # Debug: list all placeholders
                for i, ph in enumerate(slide.placeholders):
                    ph_name = getattr(ph, 'name', f'unnamed_{i}')
                    logger.info(f"  Placeholder {i}: {ph_name}")
                
                next_steps = final_data.get("next_steps", [])
                if next_steps:
                    # Try to find next_steps placeholder
                    next_steps_placeholder = (
                        find_placeholder_by_name(slide, "Subtitle 1")
                    )
                    
                    logger.info(f"Found next_steps placeholder: {next_steps_placeholder is not None}")
                    if next_steps_placeholder and next_steps_placeholder.has_text_frame:
                        # Clear any existing text first
                        next_steps_placeholder.text = ""
                        
                        # Add text properly to the text frame
                        text_frame = next_steps_placeholder.text_frame
                        text_frame.clear()
                        
                        # Add each next step as a separate paragraph with bullet
                        for i, step in enumerate(next_steps):
                            if i == 0:
                                # Use the first paragraph that already exists
                                paragraph = text_frame.paragraphs[0]
                            else:
                                # Add new paragraphs for additional steps
                                paragraph = text_frame.add_paragraph()
                            
                            paragraph.text = f"• {step}"
                            paragraph.level = 0
                        
                        # Set font size to 16 for all paragraphs in next_steps
                        try:
                            from pptx.util import Pt
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(16)
                        except Exception as e:
                            logger.warning(f"Could not set font size for next_steps: {e}")
                        
                        # Do not adjust placeholder dimensions - keep original size and position
                        logger.info(f"Set next steps with {len(next_steps)} items")
                    else:
                        logger.warning("Could not populate next_steps - no suitable placeholder found")
            else:
                logger.warning("final_slide layout not found")
        
        # Save the presentation
        prs.save(output_path)
        logger.info(f"Successfully generated PPTX presentation: {output_path}")
        
        return output_path
        
    except Exception as e:
        logger.error(f"Error generating PPTX presentation: {e}", exc_info=True)
        raise

@app.post("/make_slides/")
async def make_slides(
    doc: UploadFile = File(...),
    prompt: str = Form(...),
    slide_count: int = Form(...)
):
    # ... (Part 1: LLM call is unchanged)
    parsed_doc = await parse_doc(doc)
    doc_text = parsed_doc["text"]
    llm = ChatGoogleGenerativeAI(model="gemini-2.0-flash", temperature=0.1, api_key=GOOGLE_API_KEY)
    json_format_string = """
{
  "presentation": {
    "intro_slide": {
      "presentation_title": "The main title of the presentation in Azerbaijani",
      "presentation_date": "DD.MM.YYYY"
    },
    "project_content_slide": {
      "presentation_overview": "A 3-sentence summary of the document's content in Azerbaijani.",
      "presentation_goal": "A concise statement about the presentation's primary goal in Azerbaijani."
    },
    "content_slides": [
      {
        "general_content_title": "A title for this section of the presentation in Azerbaijani",
        "contents": [
          {
            "title": "Punkt A",
            "content-1": "A 2-3 sentence explanation for the first point in Azerbaijani."
          },
          {
            "title": "Punkt B",
            "content-2": "A 2-3 sentence explanation for the second point in Azerbaijani."
          },
          {
            "title": "Punkt C",
            "content-3": "A 2-3 sentence explanation for the third point in Azerbaijani."
          },
          {
            "title": "Punkt D",
            "content-4": "A 2-3 sentence explanation for the fourth point in Azerbaijani."
          }
        ]
      }
    ],
    "chart_slides": [
      {
        "content_title": "Maliyyə Göstəriciləri Analizi",
        "chart_summary": "Bu qrafik maliyyə göstəricilərinin il boyu dinamikasını göstərir və əsas tendensiyaları analiz edir.",
        "chart_data": [
          {
            "type": "bar",
            "title": "Aylıq Gəlir Statistikası",
            "description": "Son 6 ayın gəlir məlumatları",
            "xlabel": "Aylar",
            "ylabel": "Gəlir (min manat)",
            "x": ["Yanvar", "Fevral", "Mart", "Aprel", "May", "İyun"],
            "y": [120, 150, 180, 200, 175, 220],
            "labels": [],
            "sizes": []
          }
        ]
      },
      {
        "content_title": "Resurs Bölgüsü Analizi",
        "chart_summary": "Bu dairəvi qrafik layihə resurslarının müxtəlif kateqoriyalar üzrə bölünməsini göstərir.",
        "chart_data": [
          {
            "type": "pie",
            "title": "Layihə Resursları Paylanması",
            "description": "Resursların kateqoriyalar üzrə faiz payı",
            "xlabel": "",
            "ylabel": "",
            "x": [],
            "y": [],
            "labels": ["İnkişaf", "Test", "Dizayn", "Məlumat", "İdarəetmə"],
            "sizes": [35, 25, 20, 15, 5]
          }
        ]
      }
    ],
    "final_slide": {
      "next_steps": [
        "First next step or action item in Azerbaijani.",
        "Second next step or action item in Azerbaijani.",
        "Third next step or action item in Azerbaijani.",
        "Fourth next step or action item in Azerbaijani.",
        "Fifth next step or action item in Azerbaijani."
      ]
    }
  }
}
"""
    generation_prompt = PromptTemplate(
        input_variables=["prompt", "text", "slide_count", "json_format", "current_date"],
        template=(
            "User instruction: {prompt}\n\nDocument content:\n---\n{text}\n---\n\n"
            "**CRITICAL RULES:**\n"
            "1. `content_slides` array MUST contain EXACTLY {slide_count} objects.\n"
            "2. Inside EACH `content_slides` object, the `contents` array MUST contain EXACTLY 4 objects.\n"
            "3. The `presentation_title` in `intro_slide` MUST be a maximum of 4 words.\n"
            "4. The `next_steps` array in `final_slide` MUST contain AT LEAST 5 items.\n"
            "5. **MANDATORY**: The `chart_slides` array MUST contain EXACTLY 2 chart slide objects.\n"
            "6. **MANDATORY**: First chart MUST be type 'bar' with realistic data based on document content.\n"
            "7. **MANDATORY**: Second chart MUST be type 'pie' with realistic data based on document content.\n"
            "8. Each chart MUST have meaningful titles, labels, and data that relate to the document content.\n"
            "9. For bar charts: provide x and y arrays with at least 4-6 data points.\n"
            "10. For pie charts: provide labels and sizes arrays that sum to 100 (percentages).\n\n"
            "Return ONLY a single valid JSON object.\n\nJSON Format Example:\n{json_format}"
        )
    )
    generation_chain = generation_prompt | llm | StrOutputParser()
    current_date_str = datetime.now().strftime("%d.%m.%Y")
    llm_response = generation_chain.invoke({
        "prompt": prompt, "text": doc_text, "slide_count": slide_count,
        "json_format": json_format_string, "current_date": current_date_str
    })
    
    try:
        clean_json_str = extract_json_from_response(llm_response)
        output_json = json.loads(clean_json_str)
    except (json.JSONDecodeError, ValueError) as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate valid JSON. Error: {e}")

    # --- Part 2: PPTX generation ---
    try:
        pptx_file_path = generate_pptx_presentation(output_json)
        logger.info(f"Successfully generated PPTX presentation: {pptx_file_path}")
        
        return JSONResponse(content={
            "message": "HTML slides and PPTX presentation generated successfully.",
            "pptx_file": pptx_file_path,
            "json_data": output_json
        })
    except Exception as e:
        logger.error(f"Failed during slide generation stage: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to generate slides. Error: {e}")

@app.post("/test_pptx/")
async def test_pptx_generation():
    """
    Test endpoint to generate a sample PPTX presentation using sample data.
    Useful for testing the PPTX generation functionality.
    """
    sample_data = {
        "presentation": {
            "intro_slide": {
                "presentation_title": "Test Təqdimat",
                "presentation_date": "01.08.2025"
            },
            "project_content_slide": {
                "presentation_overview": "Bu təqdimat test məqsədilə yaradılmışdır. Bu, python-pptx kitabxanasının işlədiyini yoxlamaq üçündür.",
                "presentation_goal": "PPTX generasiya funksionallığının düzgün işlədiyini təsdiqləmək."
            },
            "content_slides": [
                {
                    "general_content_title": "Əsas Məzmun",
                    "contents": [
                        {
                            "title": "Punkt A",
                            "content-1": "Bu birinci punkt haqqında məlumatdır."
                        },
                        {
                            "title": "Punkt B", 
                            "content-2": "Bu ikinci punkt haqqında məlumatdır."
                        },
                        {
                            "title": "Punkt C",
                            "content-3": "Bu üçüncü punkt haqqında məlumatdır."
                        },
                        {
                            "title": "Punkt D",
                            "content-4": "Bu dördüncü punkt haqqında məlumatdır."
                        }
                    ]
                }
            ],
            "chart_slides": [
                {
                    "content_title": "Test Qrafik Analizi",
                    "chart_summary": "Bu qrafik test məlumatlarının dinamikasını göstərir və sistemin qrafik yaratma qabiliyyətini sınayır.",
                    "chart_data": [
                        {
                            "type": "bar",
                            "title": "Aylıq Performans Statistikası",
                            "description": "Son 6 ayın performans göstəriciləri",
                            "xlabel": "Aylar",
                            "ylabel": "Dəyər (faiz)",
                            "x": ["Yanvar", "Fevral", "Mart", "Aprel", "May", "İyun"],
                            "y": [85, 92, 78, 95, 88, 91],
                            "labels": [],
                            "sizes": []
                        }
                    ]
                },
                {
                    "content_title": "Bölgü Analizi",
                    "chart_summary": "Bu dairəvi qrafik müxtəlif kateqoriyaların paylanmasını göstərir.",
                    "chart_data": [
                        {
                            "type": "pie",
                            "title": "Resurs Bölgüsü",
                            "description": "Layihə resurslarının kateqoriyalar üzrə bölünməsi",
                            "xlabel": "",
                            "ylabel": "",
                            "x": [],
                            "y": [],
                            "labels": ["İnkişaf", "Test", "Dizayn", "Məlumat"],
                            "sizes": [40, 25, 20, 15]
                        }
                    ]
                }
            ],
            "final_slide": {
                "next_steps": [
                    "Növbəti addım: Test nəticələrini yoxlamaq",
                    "İkinci addım: Real məlumatlarla test etmək",
                    "Üçüncü addım: Nəticələri təhlil etmək",
                    "Dördüncü addım: Gələcək inkişaf planlarını müzakirə etmək",
                    "Beşinci addım: Təqdimatı təqdim etmək",
                    "Altıncı addım: İstifadəçi rəylərini toplamaq",
                    "Yeddinci addım: Təqdimatın təkmilləşdirilməsi üçün təkliflər hazırlamaq",
                    "Səkkizinci addım: Təqdimatın son versiyasını yayımlamaq",
                    "Doqquzuncu addım: Təqdimatın nəticələrini qiymətləndirmək",
                    "Onuncu addım: Gələcək layihələr üçün dərslər çıxarmaq"
                ]
            }
        }
    }
    
    try:
        pptx_file_path = generate_pptx_presentation(sample_data)
        return JSONResponse(content={
            "message": "Test PPTX presentation generated successfully.",
            "pptx_file": pptx_file_path,
            "sample_data": sample_data
        })
    except Exception as e:
        logger.error(f"Failed to generate test PPTX: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to generate test PPTX. Error: {e}")

@app.get("/")
async def serve_frontend():
    """Serve the frontend HTML page"""
    return FileResponse("static/index.html")

@app.post("/generate_json/")
async def generate_json_only(
    doc: UploadFile = File(...),
    prompt: str = Form(...),
    slide_count: int = Form(...)
):
    """Generate JSON content from document without creating PPTX"""
    global current_presentation_data
    
    # Parse document
    parsed_doc = await parse_doc(doc)
    doc_text = parsed_doc["text"]
    
    # Initialize LLM
    llm = ChatGoogleGenerativeAI(model="gemini-2.0-flash", temperature=0.1, api_key=GOOGLE_API_KEY)
    
    json_format_string = """
{
  "presentation": {
    "intro_slide": {
      "presentation_title": "The main title of the presentation in Azerbaijani",
      "presentation_date": "DD.MM.YYYY"
    },
    "project_content_slide": {
      "presentation_overview": "A 3-sentence summary of the document's content in Azerbaijani.",
      "presentation_goal": "A concise statement about the presentation's primary goal in Azerbaijani."
    },
    "content_slides": [
      {
        "general_content_title": "A title for this section of the presentation in Azerbaijani",
        "contents": [
          {
            "title": "Punkt A",
            "content-1": "A 2-3 sentence explanation for the first point in Azerbaijani."
          },
          {
            "title": "Punkt B",
            "content-2": "A 2-3 sentence explanation for the second point in Azerbaijani."
          },
          {
            "title": "Punkt C",
            "content-3": "A 2-3 sentence explanation for the third point in Azerbaijani."
          },
          {
            "title": "Punkt D",
            "content-4": "A 2-3 sentence explanation for the fourth point in Azerbaijani."
          }
        ]
      }
    ],
    "chart_slides": [
      {
        "content_title": "Maliyyə Göstəriciləri Analizi",
        "chart_summary": "Bu qrafik maliyyə göstəricilərinin il boyu dinamikasını göstərir və əsas tendensiyaları analiz edir.",
        "chart_data": [
          {
            "type": "bar",
            "title": "Aylıq Gəlir Statistikası",
            "description": "Son 6 ayın gəlir məlumatları",
            "xlabel": "Aylar",
            "ylabel": "Gəlir (min manat)",
            "x": ["Yanvar", "Fevral", "Mart", "Aprel", "May", "İyun"],
            "y": [120, 150, 180, 200, 175, 220],
            "labels": [],
            "sizes": []
          }
        ]
      },
      {
        "content_title": "Resurs Bölgüsü Analizi",
        "chart_summary": "Bu dairəvi qrafik layihə resurslarının müxtəlif kateqoriyalar üzrə bölünməsini göstərir.",
        "chart_data": [
          {
            "type": "pie",
            "title": "Layihə Resursları Paylanması",
            "description": "Resursların kateqoriyalar üzrə faiz payı",
            "xlabel": "",
            "ylabel": "",
            "x": [],
            "y": [],
            "labels": ["İnkişaf", "Test", "Dizayn", "Məlumat", "İdarəetmə"],
            "sizes": [35, 25, 20, 15, 5]
          }
        ]
      }
    ],
    "final_slide": {
      "next_steps": [
        "First next step or action item in Azerbaijani.",
        "Second next step or action item in Azerbaijani.",
        "Third next step or action item in Azerbaijani.",
        "Fourth next step or action item in Azerbaijani.",
        "Fifth next step or action item in Azerbaijani."
      ]
    }
  }
}
"""
    
    generation_prompt = PromptTemplate(
        input_variables=["prompt", "text", "slide_count", "json_format", "current_date"],
        template=(
            "User instruction: {prompt}\n\nDocument content:\n---\n{text}\n---\n\n"
            "**CRITICAL RULES:**\n"
            "1. `content_slides` array MUST contain EXACTLY {slide_count} objects.\n"
            "2. Inside EACH `content_slides` object, the `contents` array MUST contain EXACTLY 4 objects.\n"
            "3. The `presentation_title` in `intro_slide` MUST be a maximum of 4 words.\n"
            "4. The `next_steps` array in `final_slide` MUST contain AT LEAST 5 items.\n"
            "5. **MANDATORY**: The `chart_slides` array MUST contain EXACTLY 2 chart slide objects.\n"
            "6. **MANDATORY**: First chart MUST be type 'bar' with realistic data based on document content.\n"
            "7. **MANDATORY**: Second chart MUST be type 'pie' with realistic data based on document content.\n"
            "8. Each chart MUST have meaningful titles, labels, and data that relate to the document content.\n"
            "9. For bar charts: provide x and y arrays with at least 4-6 data points.\n"
            "10. For pie charts: provide labels and sizes arrays that sum to 100 (percentages).\n\n"
            "Return ONLY a single valid JSON object.\n\nJSON Format Example:\n{json_format}"
        )
    )
    
    generation_chain = generation_prompt | llm | StrOutputParser()
    current_date_str = datetime.now().strftime("%d.%m.%Y")
    
    llm_response = generation_chain.invoke({
        "prompt": prompt, "text": doc_text, "slide_count": slide_count,
        "json_format": json_format_string, "current_date": current_date_str
    })
    
    try:
        clean_json_str = extract_json_from_response(llm_response)
        output_json = json.loads(clean_json_str)
        
        # Store the generated data globally
        current_presentation_data = output_json
        
        return JSONResponse(content={
            "message": "JSON content generated successfully.",
            "json_data": output_json
        })
    except (json.JSONDecodeError, ValueError) as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate valid JSON. Error: {e}")

@app.post("/update_json/")
async def update_json(updated_data: dict = Body(...)):
    """Update the stored JSON data with user edits"""
    global current_presentation_data
    
    try:
        current_presentation_data = updated_data
        return JSONResponse(content={
            "message": "JSON data updated successfully.",
            "json_data": current_presentation_data
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to update JSON data. Error: {e}")

@app.get("/get_current_json/")
async def get_current_json():
    """Get the current stored JSON data"""
    return JSONResponse(content={
        "json_data": current_presentation_data
    })

@app.post("/generate_pptx_from_json/")
async def generate_pptx_from_json():
    """Generate PPTX presentation from the current stored JSON data"""
    global current_presentation_data
    
    if not current_presentation_data:
        raise HTTPException(status_code=400, detail="No JSON data available. Please generate content first.")
    
    try:
        pptx_file_path = generate_pptx_presentation(current_presentation_data)
        logger.info(f"Successfully generated PPTX presentation: {pptx_file_path}")
        
        return JSONResponse(content={
            "message": "PPTX presentation generated successfully.",
            "pptx_file": pptx_file_path,
            "json_data": current_presentation_data
        })
    except Exception as e:
        logger.error(f"Failed during PPTX generation: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to generate PPTX. Error: {e}")